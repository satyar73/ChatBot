import os
import json
from tqdm import tqdm
from markdownify import markdownify as md
from typing import Callable, Dict, List, Union, Optional, TypedDict, Any
from pathlib import Path
import logging
import io
from googleapiclient.http import MediaIoBaseDownload
import docx
import PyPDF2
import pptx
from langchain.docstore.document import Document
from langchain.document_loaders.base import BaseLoader

# Import from your existing project structure
import requests
from google.oauth2 import service_account
from googleapiclient.discovery import build

from app.config.chat_config import ChatConfig, chat_config
from app.services.enhancement_service import enhancement_service
from app.utils.vectorstore_client import VectorStoreClient

# Define a type for the file/folder item
class DriveItem(TypedDict):
    id: str
    name: str
    mimeType: str
    webViewLink: str
    path: str
    full_path: str
    potential_client: Optional[str]
    folder_type: Optional[str]  # one of "client", "General Resources:Case Studies", "Domain Knowledge"
    folder_name: Optional[str]  # name of the folder

# Define a type for the function return value
class FolderContents(TypedDict):
    items: List[DriveItem]
    folder_path: str


class CustomJsonLoader(BaseLoader):
    """Custom loader for JSON data"""

    def __init__(
            self, data: Union[Dict, List, Path], dataset_mapping_function: Callable[[Dict], Document]
    ):
        if isinstance(data, (dict, list)):
            self.data = data
        elif isinstance(data, (Path, str)):
            path = Path(data).resolve()
            with open(path, "r", encoding='utf-8') as f:
                self.data = json.load(f)
        self.dataset_mapping_function = dataset_mapping_function

    def load(self) -> List[Document]:
        """Load documents."""
        return list(map(self.dataset_mapping_function, self.data))

class GoogleDriveIndexer:
    """Class to handle the process of fetching Google Drive data and indexing to Pinecone"""

    def __init__(self, config: Optional[ChatConfig] = None):
        """Initialize the indexer with configuration"""
        self.config = config or chat_config
        self.last_chunks = []  # Store chunks for reporting

        # Validate configuration
        missing_settings = self.config.validate_settings()
        if missing_settings:
            raise ValueError(f"Missing required settings: {', '.join(missing_settings)}")

        # Configure logging
        logging.basicConfig(
            level=logging.INFO,
            format='%(asctime)s - %(name)s - %(levelname)s - %(message)s'
        )
        self.logger = logging.getLogger(__name__)

        # Create output directory if needed
        os.makedirs(self.config.OUTPUT_DIR, exist_ok=True)

        # Initialize Google Drive API
        self._initialize_drive_api()

    def _initialize_drive_api(self):
        """Initialize the Google Drive API client"""
        try:
            # Load service account credentials from the specified file
            credentials = service_account.Credentials.from_service_account_file(
                self.config.GOOGLE_DRIVE_CREDENTIALS_FILE,
                scopes=['https://www.googleapis.com/auth/drive.readonly']
            )

            # Build the Drive API client
            self.drive_service = build('drive', 'v3', credentials=credentials)
            self.logger.info("Successfully initialized Google Drive API client")
        except Exception as e:
            self.logger.error(f"Failed to initialize Google Drive API client: {str(e)}")
            raise

    def get_file_info(self, file_id):
        """Get the full path and name of a file or folder"""
        try:
            # First get the file/folder info
            file_info = self.drive_service.files().get(
                fileId=file_id,
                fields='id,name,parents,mimeType'
            ).execute()

            file_name = file_info.get('name', '')
            file_type = file_info.get('mimeType', '')
            parents = file_info.get('parents', [])

            # Build the path
            if not parents:
                # At root level
                path = "My Drive"
            else:
                # Get parent path
                parent_path = self.get_parent_path(parents[0])
                path = parent_path

            return {
                "id": file_id,
                "name": file_name,
                "path": path,
                "full_path": f"{path}/{file_name}",
                "type": file_type
            }

        except Exception as e:
            self.logger.error(f"Error getting file info: {str(e)}")
            return {
                "id": file_id,
                "name": "Unknown",
                "path": "Unknown",
                "full_path": f"Unknown/{file_id}",
                "type": "unknown"
            }

    def get_parent_path(self, parent_id):
        """Recursively build the path for a parent folder"""
        try:
            parent_info = self.drive_service.files().get(
                fileId=parent_id,
                fields='name,parents'
            ).execute()

            parent_name = parent_info.get('name', '')
            grandparents = parent_info.get('parents', [])

            if not grandparents:
                return f"My Drive/{parent_name}"

            # Recursively get the grandparent's path
            grandparent_path = self.get_parent_path(grandparents[0])
            return f"{grandparent_path}/{parent_name}"

        except Exception as e:
            self.logger.error(f"Error getting parent path: {str(e)}")
            return "Unknown"

    def list_folder_contents(self, folder_id=None) -> FolderContents :
        """List files and folders accessible to the service account with full paths"""
        self.logger.info(f"Listing files and folders with paths")

        # If no specific folder is requested, get all accessible files

        potential_client = "msquared"
        if folder_id is None:
            query = "trashed = false"
            folder_path = "My Drive"
            folder_name = "My Drive"
            folder_type = "Root"
        else:
            query = f"'{folder_id}' in parents and trashed = false"
            folder_info = self.get_file_info(folder_id)
            folder_name = folder_info["name"]
            folder_path = folder_info["full_path"]
            if (folder_name == "Domain Knowledge"
                or folder_name == "MSquared - Curriculum - v3.0"):
                folder_type = "Domain Knowledge"
            elif folder_name == "Case Studies":
                folder_type = "General Resources.Case Studies"
            else:
                folder_type = "client"
                potential_client = folder_name

        try:
            results = []
            page_token = None

            while True:
                response = self.drive_service.files().list(
                    q=query,
                    spaces='drive',
                    fields='nextPageToken, files(id, name, mimeType, webViewLink)',
                    pageToken=page_token,
                    pageSize=100
                ).execute()

                items = response.get('files', [])

                # Add path information to each item
                for item in items:
                    item_info = self.get_file_info(item['id'])
                    item.update({
                        'path': item_info['path'],
                        'full_path': item_info['full_path'],
                        'potential_client' : potential_client,
                        'folder_type': folder_type,
                        'folder_name': folder_name
                    })

                results.extend(items)

                page_token = response.get('nextPageToken')
                if not page_token:
                    break

            self.logger.info(f"Found {len(results)} items in {folder_path}")
            return {
                "items": results,
                "folder_path": folder_path
            }
        except Exception as e:
            self.logger.error(f"Error listing files: {str(e)}")
            return {
                "items": [],
                "folder_path": folder_path if 'folder_path' in locals() else "Unknown"
            }

    def get_supported_files(self, folder_id=None, recursive=True):
        """Get all supported files, optionally from a specific folder"""
        all_files = []
        folder_contents = self.list_folder_contents(folder_id)

        # Define supported MIME types
        supported_mime_types = [
            'application/vnd.google-apps.document',  # Google Docs
            'application/vnd.google-apps.spreadsheet',  # Google Sheets
            'application/vnd.google-apps.presentation',  # Google Slides
            'application/pdf',  # PDF files
            'application/vnd.openxmlformats-officedocument.wordprocessingml.document',  # DOCX files
            'application/vnd.openxmlformats-officedocument.presentationml.presentation',  # PPTX files
            'text/plain',  # Text files
            'text/markdown',  # Markdown files
            'text/html'  # HTML files
        ]

        for item in folder_contents["items"]:
            if item['mimeType'] in supported_mime_types:
                # It's a supported file
                all_files.append(item)
            elif (item['mimeType'] == 'application/vnd.google-apps.folder'
                  and recursive and folder_id is not None):
                # It's a folder, and we want to process recursively
                # Only process subfolders if we're starting from a specific folder
                folder_files = self.get_supported_files(item['id'], recursive=True)
                self.logger.info(f"Found {len(folder_files)} files"
                                 f"in folder {item['name']}")
                all_files.extend(folder_files)

        return all_files

    def download_and_extract_content(self, file_item):
        """Download and extract content from a Google Drive file"""
        file_id = file_item['id']
        mime_type = file_item['mimeType']
        file_name = file_item['name']

        self.logger.info(f"Processing file: {file_name} (MIME type: {mime_type})")

        try:
            # Handle Google Workspace files (need to export)
            if mime_type.startswith('application/vnd.google-apps'):
                if mime_type == 'application/vnd.google-apps.document':
                    # Export Google Docs as plain text
                    request = self.drive_service.files().export_media(fileId=file_id, mimeType='text/plain')
                    content = self._download_file(request)
                    return content.decode('utf-8')

                elif mime_type == 'application/vnd.google-apps.spreadsheet':
                    # Export Google Sheets as CSV
                    request = self.drive_service.files().export_media(fileId=file_id, mimeType='text/csv')
                    content = self._download_file(request)
                    return content.decode('utf-8')

                elif mime_type == 'application/vnd.google-apps.presentation':
                    # Check if we should use the enhanced slide extraction
                    if hasattr(self.config, 'USE_ENHANCED_SLIDES') and self.config.USE_ENHANCED_SLIDES:
                        # Use the enhanced slide extraction with GPT-4 Vision
                        # This now returns a list of slide dictionaries
                        slides = self._extract_slides_with_vision(file_id, file_name)
                        
                        # For backward compatibility, join all slides if needed
                        # The caller should check the return type and handle appropriately
                        if isinstance(slides, list):
                            # Return the first slide content for compatibility
                            # Our preprocessing step will handle the full list later
                            if slides and 'markdown' in slides[0]:
                                return slides[0]['markdown']
                            else:
                                return ""
                        else:
                            # Should not happen with the new implementation, but handle just in case
                            return slides
                    else:
                        # Fallback to basic text extraction
                        request = self.drive_service.files().export_media(fileId=file_id, mimeType='text/plain')
                        content = self._download_file(request)
                        return content.decode('utf-8')

                else:
                    self.logger.warning(f"Unsupported Google Workspace format: {mime_type}")
                    return ""

            # Handle regular files (need to download)
            else:
                request = self.drive_service.files().get_media(fileId=file_id)
                content = self._download_file(request)

                # Process based on MIME type
                if mime_type == 'application/pdf':
                    return self._extract_text_from_pdf(content)
                elif mime_type == 'application/vnd.openxmlformats-officedocument.wordprocessingml.document':
                    return self._extract_text_from_docx(content)
                elif mime_type == 'application/vnd.openxmlformats-officedocument.presentationml.presentation':
                    return self._extract_text_from_pptx(content)
                elif mime_type in ['text/plain', 'text/markdown', 'text/html']:
                    text_content = content.decode('utf-8')
                    if mime_type == 'text/html':
                        return self.html_to_markdown(text_content)
                    return text_content
                else:
                    self.logger.warning(f"Unsupported file format: {mime_type}")
                    return ""

        except Exception as e:
            self.logger.error(f"Error processing file {file_name}: {str(e)}")
            return ""

    def _download_file(self, request):
        """Download a file from Google Drive using the provided request"""
        fh = io.BytesIO()
        downloader = MediaIoBaseDownload(fh, request)
        done = False

        while not done:
            status, done = downloader.next_chunk()

        fh.seek(0)
        return fh.read()

    def _extract_text_from_pdf(self, content):
        """Extract text from a PDF file"""
        pdf_file = io.BytesIO(content)
        text = ""

        try:
            reader = PyPDF2.PdfReader(pdf_file)
            for page_num in range(len(reader.pages)):
                text += reader.pages[page_num].extract_text() + "\n\n"
            return text
        except Exception as e:
            self.logger.error(f"Error extracting text from PDF: {str(e)}")
            return ""

    def _extract_text_from_docx(self, content):
        """Extract text from a DOCX file"""
        docx_file = io.BytesIO(content)
        text = ""

        try:
            doc = docx.Document(docx_file)
            for para in doc.paragraphs:
                text += para.text + "\n"
            return text
        except Exception as e:
            self.logger.error(f"Error extracting text from DOCX: {str(e)}")
            return ""

    # Then add this new helper method to the GoogleDriveIndexer class
    def _extract_text_from_pptx(self, content):
        """Extract text from a PPTX file"""
        pptx_file = io.BytesIO(content)
        text = ""

        try:
            presentation = pptx.Presentation(pptx_file)

            # Extract title if available
            if presentation.core_properties.title:
                text += f"# {presentation.core_properties.title}\n\n"

            # Extract text from slides
            for i, slide in enumerate(presentation.slides):
                text += f"## Slide {i + 1}\n\n"

                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text += f"{shape.text.strip()}\n\n"

                # Add slide separator
                text += "---\n\n"

            return text
        except Exception as e:
            self.logger.error(f"Error extracting text from PPTX: {str(e)}")
            return ""
            
    def _extract_slides_with_vision(self, presentation_id, presentation_name):
        """
        Enhanced method to extract content from Google Slides using the Slides API
        and GPT-4 Vision for richer content extraction including visuals.
        
        Args:
            presentation_id: The ID of the Google Slides presentation
            presentation_name: The name of the presentation
            
        Returns:
            A list of dictionaries containing individual slide content and metadata
        """
        self.logger.info(f"Extracting slides with vision for {presentation_name} ({presentation_id})")
        
        try:
            # Initialize slides service if not already done
            if not hasattr(self, 'slides_service'):
                credentials = service_account.Credentials.from_service_account_file(
                    self.config.GOOGLE_DRIVE_CREDENTIALS_FILE,
                    scopes=['https://www.googleapis.com/auth/drive.readonly']
                )
                self.slides_service = build('slides', 'v1', credentials=credentials)
                
            # Get presentation
            presentation = self.slides_service.presentations().get(
                presentationId=presentation_id
            ).execute()
            
            # Get all slides
            slides = presentation.get('slides', [])
            self.logger.info(f"Processing {len(slides)} slides from {presentation_name}")
            
            # Store individual slide content
            slide_contents = []
            
            # Add presentation info as the first item
            presentation_info = {
                'title': f"{presentation_name} - Overview",
                'markdown': f"# {presentation_name}\n\nThis presentation contains {len(slides)} slides.",
                'slide_number': 0,
                'is_slide': True,
                'slide_title': presentation_name
            }
            slide_contents.append(presentation_info)
            
            # Limit number of slides to process if specified in config
            max_slides = getattr(self.config, 'MAX_SLIDES_TO_PROCESS', len(slides))
            slides_to_process = slides[:max_slides]
            
            # Process each slide
            for slide_index, slide in enumerate(slides_to_process):
                slide_number = slide_index + 1
                self.logger.info(f"Processing slide {slide_number}/{len(slides)}")
                
                # Export the slide as PNG
                try:
                    export_request = self.slides_service.presentations().pages().getThumbnail(
                        presentationId=presentation_id,
                        pageObjectId=slide['objectId'],
                        thumbnailProperties_thumbnailSize='LARGE'
                    ).execute()
                    
                    # Get the thumbnail URL
                    thumbnail_url = export_request.get('contentUrl')
                    
                    # Download the image
                    image_response = requests.get(thumbnail_url, timeout=30)
                    if image_response.status_code != 200:
                        self.logger.warning(f"Failed to download slide {slide_number} image: {image_response.status_code}")
                        continue
                        
                    image_content = image_response.content
                    
                    # Analyze with GPT-4 Vision
                    slide_markdown = self._analyze_slide_with_llm(image_content, slide_number)
                    if slide_markdown:
                        # Try to extract title from the markdown if it exists
                        slide_title = None
                        if slide_markdown.startswith('## '):
                            title_line = slide_markdown.split('\n')[0]
                            slide_title = title_line.replace('## ', '').strip()
                        
                        # If no title found, try to derive one from the content
                        if not slide_title:
                            # Look for text elements that might be titles
                            for element in slide.get('pageElements', []):
                                if 'shape' in element and 'text' in element['shape']:
                                    text_content = ""
                                    for textElement in element['shape']['text'].get('textElements', []):
                                        if 'textRun' in textElement and 'content' in textElement['textRun']:
                                            text_content += textElement['textRun']['content']
                                    
                                    # Check if this might be a title based on position or formatting
                                    if ('title' in element.get('objectId', '').lower() or
                                        ('transform' in element and element['transform'].get('scaleY', 0) > 1) or
                                        (len(text_content.strip()) < 100 and text_content.strip())):  # Short text is likely a title
                                        slide_title = text_content.strip()
                                        break
                        
                        # If still no title, use a generic one
                        if not slide_title:
                            slide_title = "Untitled Slide"
                        
                        # Format the content with slide number and title
                        formatted_content = f"## Slide {slide_number}\n\n"
                        formatted_content += f"### {slide_title}\n\n"
                        formatted_content += slide_markdown
                        
                        # Store as individual slide
                        slide_contents.append({
                            'title': f"{presentation_name} - Slide {slide_number}",
                            'markdown': formatted_content,
                            'slide_number': slide_number,
                            'is_slide': True,
                            'slide_title': slide_title
                        })
                    else:
                        # Fallback to basic extraction if vision analysis fails
                        slide_title = None
                        slide_content = []
                        
                        # Extract text elements
                        for element in slide.get('pageElements', []):
                            if 'shape' in element and 'text' in element['shape']:
                                text_content = ""
                                for textElement in element['shape']['text'].get('textElements', []):
                                    if 'textRun' in textElement and 'content' in textElement['textRun']:
                                        text_content += textElement['textRun']['content']
                                
                                # Check if this might be a title
                                if ('title' in element.get('objectId', '').lower() or
                                    ('transform' in element and element['transform'].get('scaleY', 0) > 1) or
                                    (len(text_content.strip()) < 100 and text_content.strip())):  # Short text is likely a title
                                    slide_title = text_content.strip()
                                else:
                                    # Add to content if not empty
                                    if text_content.strip():
                                        slide_content.append(text_content.strip())
                        
                        # If no title found, use a generic one
                        if not slide_title:
                            slide_title = "Untitled Slide"
                        
                        # Format as markdown with slide number and title
                        formatted_content = f"## Slide {slide_number}\n\n"
                        formatted_content += f"### {slide_title}\n\n"
                        for item in slide_content:
                            formatted_content += f"{item}\n\n"
                            
                        # Store as individual slide
                        slide_contents.append({
                            'title': f"{presentation_name} - Slide {slide_number}",
                            'markdown': formatted_content,
                            'slide_number': slide_number,
                            'is_slide': True,
                            'slide_title': slide_title
                        })
                
                except Exception as e:
                    self.logger.error(f"Error processing slide {slide_number}: {str(e)}")
                    continue
            
            return slide_contents
            
        except Exception as e:
            self.logger.error(f"Error extracting slides with vision: {str(e)}")
            # Fall back to plain text extraction and let the caller handle it
            request = self.drive_service.files().export_media(fileId=presentation_id, mimeType='text/plain')
            content = self._download_file(request)
            return [{
                'title': presentation_name,
                'markdown': content.decode('utf-8'),
                'is_slide': False
            }]
    
    def _analyze_slide_with_llm(self, image_content, slide_number):
        """
        Analyze a slide image using GPT-4 Vision
        
        Args:
            image_content: Raw image data
            slide_number: The slide number for reference
            
        Returns:
            Markdown string with the analysis result
        """
        prompt = self.config.GOOGLE_SLIDE_IMAGE_TO_TEXT_PROMPT
        # Use the enhancement service to analyze the image
        return enhancement_service.analyze_image_with_llm(image_content, prompt, getattr(self.config, 'OPENAI_VISION_MODEL', 'gpt-4o'))

    def condense_content_using_llm(self, content):
        """Summarize content using OpenAI's API"""
        # Use the enhancement service to condense content
        return enhancement_service.condense_content_using_llm(content)

    def html_to_markdown(self, html_text):
        """Convert HTML to Markdown format"""
        markdown_text = md(html_text, newline_style="BACKSLASH", default_title=True, heading_style='ATX')
        markdown_text = markdown_text.replace('\n\n', '\n')

        if self.config.SUMMARIZE_CONTENT:
            self.logger.info("Summarizing content using LLM...")
            markdown_text = self.condense_content_using_llm(markdown_text)

        return markdown_text

    def _preprocess_slide_content(self, content: str, file_name: str) -> List[Dict[str, Any]]:
        """
        Preprocess slide content to preserve slide boundaries.
        Splits the content at slide boundaries (marked by "## Slide X" or "---") and
        creates separate records for each slide.
        
        Args:
            content: The slide content as markdown
            file_name: The name of the slide presentation
            
        Returns:
            List of records, one for each slide
        """
        self.logger.info(f"Preprocessing slides for {file_name}")
        
        # Check if this is slide content (has slide markers)
        if "## Slide " not in content and "---" not in content:
            # Not slide content or no clear slide boundaries
            return [{
                'title': file_name,
                'markdown': content,
                'is_slide': False
            }]
            
        # Split the content at slide boundaries
        # We can use either "## Slide X" markers or "---" horizontal rule markers
        slide_records = []
        
        # First look for slide markers
        if "## Slide " in content:
            # Split by slide headers
            parts = content.split("## Slide ")
            
            # The first part might be presentation title or intro
            if parts[0].strip():
                intro = parts[0].strip()
                slide_records.append({
                    'title': f"{file_name} - Introduction",
                    'markdown': intro,
                    'is_slide': True,
                    'slide_number': 0
                })
            
            # Process each slide part
            for i, part in enumerate(parts[1:], 1):
                # Split at first newline to separate slide number from content
                slide_parts = part.split('\n', 1)
                if len(slide_parts) > 1:
                    slide_number = slide_parts[0].strip()
                    slide_content = slide_parts[1].strip()
                    
                    # Clean up the content by removing horizontal rules at the end
                    if slide_content.endswith("---"):
                        slide_content = slide_content.rsplit("---", 1)[0].strip()
                    
                    slide_records.append({
                        'title': f"{file_name} - Slide {slide_number}",
                        'markdown': f"## Slide {slide_number}\n\n{slide_content}",
                        'is_slide': True,
                        'slide_number': i
                    })
        else:
            # Use horizontal rules as slide separators
            parts = content.split("---")
            
            for i, part in enumerate(parts):
                part = part.strip()
                if part:
                    slide_records.append({
                        'title': f"{file_name} - Slide {i+1}",
                        'markdown': part,
                        'is_slide': True,
                        'slide_number': i+1
                    })
        
        self.logger.info(f"Split presentation into {len(slide_records)} slides")
        return slide_records

    def prepare_drive_documents(self):
        """Process all supported files from Google Drive"""
        self.logger.info("Fetching and processing Google Drive documents...")

        # Get folder_id from config if available, otherwise use None (root)
        folder_id = getattr(self.config, 'GOOGLE_DRIVE_FOLDER_ID', None)
        recursive = getattr(self.config, 'GOOGLE_DRIVE_RECURSIVE', True)

        files = self.get_supported_files(folder_id, recursive)
        records = []

        for file in tqdm(files):
            try:
                self.logger.info(f"Processing {file['name']}...")

                # Check if this is a presentation with enhanced vision processing
                is_presentation = (file['mimeType'] == 'application/vnd.google-apps.presentation' or 
                                  file['mimeType'] == 'application/vnd.openxmlformats-officedocument.presentationml.presentation')
                use_enhanced_slides = hasattr(self.config, 'USE_ENHANCED_SLIDES') and self.config.USE_ENHANCED_SLIDES
                
                # Special handling for presentations with enhanced vision extraction
                if is_presentation and use_enhanced_slides:
                    try:
                        # Extract slide content using vision
                        slide_contents = self._extract_slides_with_vision(file['id'], file['name'])
                        
                        # Process each slide individually
                        if isinstance(slide_contents, list) and len(slide_contents) > 0:
                            self.logger.info(f"Processing {len(slide_contents)} slides from vision extraction")
                            
                            for slide_content in slide_contents:
                                # Create record for each slide
                                record = {
                                    'title': slide_content['title'],
                                    'url': file.get('webViewLink', ''),
                                    'source': 'Google Drive',
                                    'type': file['folder_type'],  # Preserve original folder type (e.g., "client")
                                    'document_type': 'presentation_slide',  # Add document type as additional metadata
                                    'client': file['potential_client'],
                                    'markdown': slide_content['markdown'],
                                    'parent_presentation': file['name'],
                                    'slide_number': slide_content.get('slide_number', 0)
                                }
                                records.append(record)
                            
                            # Skip the regular content processing
                            continue
                    except Exception as e:
                        self.logger.error(f"Error processing enhanced slides, falling back to standard processing: {str(e)}")
                
                # Standard processing for non-presentations or if enhanced processing failed
                content = self.download_and_extract_content(file)

                # Skip if content extraction failed
                if not content.strip():
                    self.logger.warning(f"Empty content for file: {file['name']}")
                    continue

                # Add file title if not in content
                if file['name'] not in content:
                    content = f"# {file['name']}\n\n{content}"

                # Summarize if needed
                if self.config.SUMMARIZE_CONTENT:
                    content = self.condense_content_using_llm(content)
                
                # Standard presentation processing (without enhanced vision)
                if is_presentation and not use_enhanced_slides:
                    # Process slides individually to preserve boundaries
                    slide_records = self._preprocess_slide_content(content, file['name'])
                    
                    for slide_record in slide_records:
                        # Create record for each slide
                        record = {
                            'title': slide_record['title'],
                            'url': file.get('webViewLink', ''),
                            'source': 'Google Drive',
                            'type': file['folder_type'],  # Preserve original folder type (e.g., "client")
                            'document_type': 'presentation_slide',  # Add document type as additional metadata
                            'client': file['potential_client'],
                            'markdown': slide_record['markdown'],
                            'parent_presentation': file['name'],
                            'slide_number': slide_record.get('slide_number', 0)
                        }
                        records.append(record)
                else:
                    # Regular document processing
                    record = {
                        'title': file['name'],
                        'url': file.get('webViewLink', ''),
                        'source': 'Google Drive',
                        'type': file['folder_type'],
                        'client': file['potential_client'],
                        'markdown': content
                    }
                    records.append(record)

            except Exception as ex:
                self.logger.error(f"Could not process {file['name']} due to: {ex}")

        self.logger.info(f"Prepared {len(records)} document records!")

        # Save intermediate file if needed
        if self.config.SAVE_INTERMEDIATE_FILES:
            processed_file = os.path.join(self.config.OUTPUT_DIR, "drive_processed.json")
            with open(processed_file, 'w') as f:
                json.dump(records, f, indent=2)

        return records

    def get_embedding_dimensions(self, model_name: str) -> int:
        """Get the embedding dimensions for a given model"""
        model_dimensions = {
            "text-embedding-3-small": 1536,
            "text-embedding-3-large": 3072,
            "text-embedding-ada-002": 1536
        }
        return model_dimensions.get(model_name, 1536)  # Default to 1536 if unknown

    def _get_file_type(self, mime_type: str) -> str:
        """
        Get a human-readable file type from MIME type.
        
        Args:
            mime_type: MIME type of the file
            
        Returns:
            Human-readable file type
        """
        mime_to_type = {
            'application/vnd.google-apps.presentation': 'presentation',
            'application/vnd.google-apps.document': 'document',
            'application/vnd.google-apps.spreadsheet': 'spreadsheet',
            'application/pdf': 'pdf',
            'image/jpeg': 'image',
            'image/png': 'image',
            'text/plain': 'text'
        }
        return mime_to_type.get(mime_type, 'unknown')

    def process_drive(self, folder_id: str = None) -> List[Dict[str, Any]]:
        """Process all supported files from Google Drive"""
        self.logger.info("Fetching and processing Google Drive documents...")

        # Get folder_id from config if available, otherwise use None (root)
        folder_id = folder_id or getattr(self.config, 'GOOGLE_DRIVE_FOLDER_ID', None)
        recursive = getattr(self.config, 'GOOGLE_DRIVE_RECURSIVE', True)

        files = self.get_supported_files(folder_id, recursive)
        records = []

        for file in tqdm(files):
            try:
                self.logger.info(f"Processing {file['name']}...")

                # Download and extract content
                content = self.download_and_extract_content(file)

                # Skip if content extraction failed
                if not content.strip():
                    self.logger.warning(f"Empty content for file: {file['name']}")
                    continue

                # Add file title if not in content
                if file['name'] not in content:
                    content = f"# {file['name']}\n\n{content}"

                # Summarize if needed
                if self.config.SUMMARIZE_CONTENT:
                    content = self.condense_content_using_llm(content)

                # Create record
                record = {
                    'title': file['name'],
                    'url': file.get('webViewLink', ''),
                    'source': 'Google Drive',
                    'type': self._get_file_type(file['mimeType']),
                    'client': file['potential_client'],
                    'markdown': content
                }
                records.append(record)

            except Exception as ex:
                self.logger.error(f"Could not process {file['name']} due to: {ex}")

        self.logger.info(f"Prepared {len(records)} document records!")

        # Save intermediate file if needed
        if self.config.SAVE_INTERMEDIATE_FILES:
            processed_file = os.path.join(self.config.OUTPUT_DIR, "drive_processed.json")
            with open(processed_file, 'w') as f:
                json.dump(records, f, indent=2)

        return records

    def run_full_process(self):
        """Initialize the Google Drive API and validate settings"""
        try:
            # Validate credentials and API access
            files = self.list_folder_contents()
            
            self.logger.info(f"Successfully connected to Google Drive API, found {len(files)} accessible files")
            return {
                "status": "success", 
                "message": f"Successfully connected to Google Drive API, found {len(files)} accessible files",
                "files_count": len(files)
            }

        except Exception as e:
            self.logger.error(f"Error initializing Google Drive API: {str(e)}")
            return {"status": "error", "message": str(e)}

    def get_google_drive_files(self) -> Dict[str, Any]:
        """Get list of indexed Google Drive files"""
        try:
            # Try to load the Google Drive processed files
            drive_path = os.path.join(self.config.OUTPUT_DIR, "drive_processed.json")
            
            if os.path.exists(drive_path):
                with open(drive_path, "r") as f:
                    files = json.load(f)
                    
                # Extract basic file information
                file_list = [
                    {
                        "id": idx,
                        "title": file.get("title", "Unknown"),
                        "url": file.get("url", ""),
                        "size": len(file.get("markdown", "")) if "markdown" in file else 0
                    }
                    for idx, file in enumerate(files)
                ]
                
                return {
                    "status": "success",
                    "files": file_list,
                    "count": len(file_list)
                }
            else:
                # Check if we can query the vector store directly
                try:
                    total_vector_count = 0
                    for chat_model_config in chat_config.chat_model_configs.values():
                        vector_store_config = chat_model_config.vector_store_config
                        vector_store_client: VectorStoreClient = VectorStoreClient.get_vector_store_client(vector_store_config)
                        total_vector_count += vector_store_client.get_vector_count()
                    
                    #TODO we should not sum the vector count across two different indexes; UI need to pass the index_name
                    if total_vector_count > 0:
                        return {
                            "status": "success",
                            "files": [],
                            "count": 0,
                            "vector_count": total_vector_count,
                            "message": "Drive file list not available, but vectors are in the index"
                        }
                except Exception as e:
                    self.logger.error(f"Error querying Pinecone for Google Drive files: {str(e)}")
                
                # No processed files available
                return {
                    "status": "success",
                    "files": [],
                    "count": 0,
                    "message": "No Google Drive files indexed or file list not available"
                }
                
        except Exception as e:
            self.logger.error(f"Error retrieving Google Drive files: {str(e)}")
            return {"status": "error", "message": str(e)}